import collections.abc
import collections
collections.Mapping = collections.abc.Mapping
collections.Sequence = collections.abc.Sequence

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw

# --- Configuration ---
BASE_DIR = r"c:\Users\Hardik Verma\Downloads\Ageon-20260414T094340Z-3-001"
PPT_OUTPUT = os.path.join(BASE_DIR, "ageon_presentation.pptx")
LOGO_PATH = os.path.join(BASE_DIR, "Ageon - logo.jpg")

# --- Brand Colors (from PDF / brand book) ---
TEAL       = RGBColor(0, 168, 150)     # Primary teal from logo
DARK_TEAL  = RGBColor(10, 52, 48)      # Deep dark accent
RED        = RGBColor(225, 37, 29)     # Red from logo mark
GOLD       = RGBColor(248, 177, 40)    # Gold/yellow from logo mark
BLACK_TEXT = RGBColor(30, 30, 30)      # Near-black for body copy
GREY_TEXT  = RGBColor(80, 80, 80)      # Lighter body copy
WHITE      = RGBColor(255, 255, 255)

# --- Exact content from AGEON Concept.docx (verbatim) ---

# ============ SLIDE 1: WHAT AGEON IS & WHAT AGEON DOES ============

SLIDE1_MAIN_TITLE = "WHAT AGEON IS & WHAT AGEON DOES"

SLIDE1_SECTION_A_TITLE = "What AGEON Is"
SLIDE1_SECTION_A_BODY = (
    "AGEON is India\u2019s Smart-Ageing & Longevity Platform\u2014a category-defining system "
    "built to transform how people experience health, ageing, and performance.\n"
    "It is not a wellness clinic and not a collection of therapies.\n"
    "AGEON is a protocol-driven, science-backed longevity system designed for "
    "measurable health outcomes and long-term vitality."
)

SLIDE1_WHY_TITLE = "AGEON exists because:"
SLIDE1_WHY_BODY = (
    "India is living longer, but not healthier.\n"
    "Modern lifestyles are accelerating:\n"
    "\u2022 Fatigue & burnout\n"
    "\u2022 Metabolic dysfunction\n"
    "\u2022 Poor recovery\n"
    "\u2022 Hormonal imbalance\n"
    "\u2022 Early ageing\n\n"
    "Most people act only after symptoms appear.\n"
    "AGEON shifts this paradigm to preventive, proactive, and performance-led health."
)

SLIDE1_SECTION_B_TITLE = "What AGEON Does"
SLIDE1_SECTION_B_INTRO = "AGEON integrates diagnostics, therapies, protocols, and tracking into one seamless ecosystem."

SLIDE1_CARDS = [
    (
        "1. Protocol-Led Longevity System",
        "AGEON does not sell random treatments.\nIt delivers structured, outcome-driven protocols."
    ),
    (
        "2. Integrated Diagnostics Engine",
        "AGEON begins with deep health measurement, not assumptions.\n"
        "This creates a true biological baseline, enabling precision-driven care."
    ),
    (
        "3. Therapy Stack (Delivered via Protocols)",
        "Therapies are not sold in isolation, but assigned based on need."
    ),
    (
        "4. The AGEON Journey (Structured System)",
        "Every member follows a defined transformation pathway:\n"
        "Assessment \u2192 Full diagnostic baseline\n"
        "Protocol Assignment \u2192 Personalized roadmap\n"
        "Guided Delivery \u2192 Expert-led therapies\n"
        "Tracking & Optimization \u2192 Continuous improvement\n"
        "This ensures measurable outcomes, not guesswork."
    ),
    (
        "5. Membership-First Model",
        "AGEON is built on continuity, not one-time visits.\n"
        "This makes AGEON a lifestyle system\u2014not a service center."
    ),
]

# ============ SLIDE 2: PURPOSE, VISION & SCIENCE BEHIND AGEON ============

SLIDE2_MAIN_TITLE = "PURPOSE, VISION & SCIENCE BEHIND AGEON"

SLIDE2_PURPOSE_TITLE = "The Purpose"
SLIDE2_PURPOSE_BODY = (
    "AGEON\u2019s purpose is clear:\n"
    "To shift healthcare from reactive treatment to proactive longevity.\n\n"
    "It focuses on:\n"
    "\u2022 Extending healthspan (not just lifespan)\n"
    "\u2022 Preventing disease before it appears\n"
    "\u2022 Enhancing physical, mental, and metabolic performance"
)

SLIDE2_VISION_TITLE = "The Vision"
SLIDE2_VISION_BODY = (
    "AGEON is building:\nIndia\u2019s Longevity Infrastructure\n"
    "Not a clinic chain. Not a therapy brand.\n"
    "But a national platform for smart-ageing and preventive health.\n\n"
    "The long-term vision includes:\n"
    "\u2022 Multi-city expansion (flagship \u2192 metro \u2192 tier 2 \u2192 corporate)\n"
    "\u2022 Standardized, scalable wellness systems\n"
    "\u2022 A premium, trust-driven health ecosystem"
)

SLIDE2_SCIENCE_TITLE = "The Science Behind AGEON"
SLIDE2_SCIENCE_INTRO = "AGEON is built on a data-driven, cellular-level health model."

SLIDE2_SCIENCE_CARDS = [
    (
        "1. Integrated Health Measurement",
        "AGEON tracks key biological systems:\n"
        "\u2022 Metabolic health (energy, glucose, body composition)\n"
        "\u2022 Cardiovascular risk (endurance, early indicators)\n"
        "\u2022 Biological ageing (cellular decline, vitality markers)\n"
        "\u2022 Inflammation & recovery (stress load, resilience)\n"
        "This creates a 360\u00b0 view of human health, not fragmented insights."
    ),
    (
        "2. Cellular & Functional Optimization",
        "AGEON targets root causes of ageing:\n"
        "\u2022 Cellular damage\n\u2022 Oxygen inefficiency\n"
        "\u2022 Chronic inflammation\n\u2022 Metabolic imbalance\n\n"
        "Therapies are designed to:\n"
        "\u2022 Improve cellular repair\n"
        "\u2022 Enhance mitochondrial efficiency\n"
        "\u2022 Boost recovery and resilience"
    ),
    (
        "3. Protocol-Based Clinical Logic",
        "Unlike traditional wellness:\n"
        "\u2022 No random treatments\n\u2022 No session-based approach\n\n"
        "AGEON follows:\n"
        "\u2022 Doctor-led intake\n\u2022 Data-driven decision systems\n"
        "\u2022 Standardized clinical pathways\n\u2022 Continuous outcome tracking"
    ),
    (
        "4. Measurable Outcomes & Optimization",
        "AGEON ensures:\n"
        "\u2022 Continuous tracking\n\u2022 Protocol refinement\n"
        "\u2022 Long-term performance visibility\n\n"
        "This transforms wellness into a measurable, repeatable system."
    ),
]

SLIDE2_CLOSING_TITLE = "The Ultimate Goal"
SLIDE2_CLOSING_BODY = (
    "To make smart-ageing a mainstream lifestyle in India.\n\n"
    "AGEON is building a future where:\n"
    "\u2022 Ageing is optimized\n"
    "\u2022 Performance is sustained\n"
    "\u2022 Health is engineered, not managed\n\n"
    "AGEON is not a clinic.\n"
    "It is not a therapy room. It is a category-defining smart-ageing platform."
)


# ─── Helper functions ───────────────────────────────────────────

def make_light_bg(filename):
    """Create a clean, white background with very subtle teal accents."""
    w, h = 1920, 1080
    img = Image.new('RGB', (w, h), (255, 255, 255))
    draw = ImageDraw.Draw(img)

    # Subtle warm-white teal gradient orbs
    draw.ellipse((-200, -200, 600, 600), fill=(240, 252, 250))
    draw.ellipse((w - 500, h - 400, w + 300, h + 300), fill=(235, 249, 247))

    # Very faint teal accent line at bottom
    draw.rectangle([(0, h - 6), (w, h)], fill=(0, 168, 150))

    path = os.path.join(BASE_DIR, filename)
    img.save(path, quality=95)
    return path


def add_text(slide, left, top, width, height, text,
             size=Pt(12), bold=False, color=BLACK_TEXT,
             align=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add an editable text box with word wrap."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tb


def add_card(slide, left, top, width, height,
             title_text, body_text,
             title_color=TEAL, body_color=GREY_TEXT):
    """Draw a teal-bordered card with title and body text."""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(200, 230, 226)  # very light teal border
    card.line.width = Pt(1.2)

    # Title inside card
    add_text(slide,
             left + Inches(0.15), top + Inches(0.08),
             width - Inches(0.3), Inches(0.35),
             title_text, size=Pt(11), bold=True, color=title_color)

    # Body inside card
    add_text(slide,
             left + Inches(0.15), top + Inches(0.38),
             width - Inches(0.3), height - Inches(0.45),
             body_text, size=Pt(9), color=body_color)


# ─── Build Presentation ─────────────────────────────────────────

def build():
    if os.path.exists(PPT_OUTPUT):
        try:
            os.remove(PPT_OUTPUT)
        except PermissionError:
            print(f"ERROR: Close {PPT_OUTPUT} in PowerPoint first.")
            return

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    bg_path = make_light_bg("bg_light.png")

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 1
    # ══════════════════════════════════════════════════════════════
    s1 = prs.slides.add_slide(blank)
    s1.shapes.add_picture(bg_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))

    # Logo (top-left)
    if os.path.exists(LOGO_PATH):
        s1.shapes.add_picture(LOGO_PATH, Inches(0.6), Inches(0.35), width=Inches(2.8))

    # Main title bar
    add_text(s1, Inches(0.6), Inches(1.4), Inches(12), Inches(0.6),
             SLIDE1_MAIN_TITLE,
             size=Pt(28), bold=True, color=DARK_TEAL, font_name="Montserrat")

    # Teal accent line under title
    line = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.6), Inches(2.0), Inches(2), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = TEAL
    line.line.fill.background()

    # --- Left column: What AGEON Is + Why ---
    col_left = Inches(0.6)
    col_w = Inches(5.2)

    add_text(s1, col_left, Inches(2.2), col_w, Inches(0.4),
             SLIDE1_SECTION_A_TITLE,
             size=Pt(18), bold=True, color=TEAL)

    add_text(s1, col_left, Inches(2.65), col_w, Inches(1.6),
             SLIDE1_SECTION_A_BODY,
             size=Pt(10), color=GREY_TEXT)

    add_text(s1, col_left, Inches(4.3), col_w, Inches(0.35),
             SLIDE1_WHY_TITLE,
             size=Pt(12), bold=True, color=RED)

    add_text(s1, col_left, Inches(4.65), col_w, Inches(2.5),
             SLIDE1_WHY_BODY,
             size=Pt(9), color=GREY_TEXT)

    # --- Right column: What AGEON Does (cards) ---
    right_x = Inches(6.2)
    right_w = Inches(6.8)

    add_text(s1, right_x, Inches(2.2), right_w, Inches(0.4),
             SLIDE1_SECTION_B_TITLE,
             size=Pt(18), bold=True, color=TEAL)

    add_text(s1, right_x, Inches(2.65), right_w, Inches(0.35),
             SLIDE1_SECTION_B_INTRO,
             size=Pt(10), color=GREY_TEXT)

    # 5 cards in a 2-col grid
    card_w = Inches(3.25)
    card_h = Inches(1.35)
    gap_x = Inches(0.2)
    gap_y = Inches(0.15)
    start_y = Inches(3.1)

    positions = [
        (right_x,                start_y),
        (right_x + card_w + gap_x, start_y),
        (right_x,                start_y + card_h + gap_y),
        (right_x + card_w + gap_x, start_y + card_h + gap_y),
        (right_x,                start_y + 2 * (card_h + gap_y)),
    ]

    for i, (title, body) in enumerate(SLIDE1_CARDS):
        x, y = positions[i]
        add_card(s1, x, y, card_w, card_h, title, body)

    # ══════════════════════════════════════════════════════════════
    #  SLIDE 2
    # ══════════════════════════════════════════════════════════════
    s2 = prs.slides.add_slide(blank)
    s2.shapes.add_picture(bg_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))

    # Logo
    if os.path.exists(LOGO_PATH):
        s2.shapes.add_picture(LOGO_PATH, Inches(0.6), Inches(0.35), width=Inches(2.8))

    # Main title
    add_text(s2, Inches(0.6), Inches(1.4), Inches(12), Inches(0.6),
             SLIDE2_MAIN_TITLE,
             size=Pt(28), bold=True, color=DARK_TEAL, font_name="Montserrat")

    line2 = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0.6), Inches(2.0), Inches(2), Inches(0.04))
    line2.fill.solid()
    line2.fill.fore_color.rgb = TEAL
    line2.line.fill.background()

    # --- Top row: Purpose (left) + Vision (right) ---
    box_w = Inches(6.0)
    box_h = Inches(2.2)
    row1_y = Inches(2.2)

    # Purpose
    add_text(s2, Inches(0.6), row1_y, box_w, Inches(0.35),
             SLIDE2_PURPOSE_TITLE, size=Pt(16), bold=True, color=TEAL)
    add_text(s2, Inches(0.6), row1_y + Inches(0.4), box_w, box_h - Inches(0.4),
             SLIDE2_PURPOSE_BODY, size=Pt(10), color=GREY_TEXT)

    # Vision
    add_text(s2, Inches(6.9), row1_y, box_w, Inches(0.35),
             SLIDE2_VISION_TITLE, size=Pt(16), bold=True, color=TEAL)
    add_text(s2, Inches(6.9), row1_y + Inches(0.4), box_w, box_h - Inches(0.4),
             SLIDE2_VISION_BODY, size=Pt(10), color=GREY_TEXT)

    # --- Science section title ---
    sci_y = Inches(4.5)
    add_text(s2, Inches(0.6), sci_y, Inches(6), Inches(0.35),
             SLIDE2_SCIENCE_TITLE, size=Pt(16), bold=True, color=TEAL)
    add_text(s2, Inches(0.6), sci_y + Inches(0.35), Inches(8), Inches(0.25),
             SLIDE2_SCIENCE_INTRO, size=Pt(10), color=GREY_TEXT)

    # 4 science cards in a row
    sci_card_w = Inches(3.0)
    sci_card_h = Inches(2.2)
    sci_start_x = Inches(0.6)
    sci_start_y = sci_y + Inches(0.7)
    sci_gap = Inches(0.15)

    for i, (title, body) in enumerate(SLIDE2_SCIENCE_CARDS):
        cx = sci_start_x + i * (sci_card_w + sci_gap)
        add_card(s2, cx, sci_start_y, sci_card_w, sci_card_h,
                 title, body, title_color=DARK_TEAL, body_color=GREY_TEXT)

    prs.save(PPT_OUTPUT)
    print(f"Successfully saved: {PPT_OUTPUT}")


if __name__ == "__main__":
    build()
